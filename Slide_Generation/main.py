from fastapi import FastAPI, File, UploadFile
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
import os
from io import BytesIO
import fitz
from transformers import pipeline
from pptx import Presentation

app = FastAPI()

allowed_origins = os.getenv("ALLOWED_ORIGINS", "http://localhost:3000").split(",")
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_credentials=True,
    allow_methods=["GET", "POST"],
    allow_headers=["*"],
)

summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

def extract_text_from_pdf(pdf_path):
    document = fitz.open(pdf_path)
    full_text = ""
    for page_num in range(len(document)):
        page = document.load_page(page_num)
        full_text += page.get_text()
    return full_text

def summarize_text(text, max_length=150, chunk_size=1024):
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    summaries = []
    for chunk in chunks:
        summary = summarizer(chunk, max_length=max_length, min_length=50, do_sample=False)
        summaries.append(summary[0]['summary_text'])
    return summaries

def extract_topic(text):
    sentences = text.split('. ')
    topic = '. '.join(sentences[:2])
    return topic

def create_presentation(topic, summaries):
    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = "Lecture Topic:"
    subtitle_placeholder.text = topic

    for summary in summaries:
        bullet_points = summary.split('. ')
        for i in range(0, len(bullet_points), 2):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            body_placeholder = slide.shapes.placeholders[1]
            if i == 0:
                title_placeholder.text = "Key Points"
            body_placeholder.text = "\n".join([f"{point.strip()}" for point in bullet_points[i:i+2]])

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.shapes.placeholders[1]
    title_placeholder.text = "Summary"
    body_placeholder.text = summaries[-1]

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

@app.post("/generate_presentation/")
async def generate_presentation(file: UploadFile = File(...)):
    pdf_text = await file.read()
    pdf_path = "temp.pdf"
    with open(pdf_path, "wb") as f:
        f.write(pdf_text)

    print("Extracting test from PDF...")
    pdf_text = extract_text_from_pdf(pdf_path)
    print("Summarizing text...")
    summaries = summarize_text(pdf_text)
    print("Extracting topic...")
    topic = extract_topic(pdf_text)
    print("Creating presentation...")
    ppt_stream = create_presentation(topic, summaries)

    ppt_filename = "generated_presentation.pptx"
    with open(ppt_filename, "wb") as f:
        f.write(ppt_stream.read())

    return FileResponse(ppt_filename, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename=ppt_filename)