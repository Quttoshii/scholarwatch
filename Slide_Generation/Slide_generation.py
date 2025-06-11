import fitz  # PyMuPDF for PDF extraction
from transformers import pipeline  # For summarization
from pptx import Presentation  # For creating PowerPoint
from io import BytesIO

# Initialize summarization pipeline
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# Step 1: Extract text from PDF
def extract_text_from_pdf(pdf_path):
    document = fitz.open(pdf_path)
    full_text = ""
    for page_num in range(len(document)):
        page = document.load_page(page_num)
        full_text += page.get_text()
    return full_text

# Step 2: Summarize the extracted text
def summarize_text(text, max_length=150, chunk_size=1024):
    # Split the text into chunks of 1024 characters
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    
    # Summarize each chunk and collect the summaries
    summaries = []
    for chunk in chunks:
        summary = summarizer(chunk, max_length=max_length, min_length=50, do_sample=False)
        summaries.append(summary[0]['summary_text'])
    
    return summaries

# Step 3: Extract the topic from the first few sentences of the PDF
def extract_topic(text):
    sentences = text.split('. ')
    topic = '. '.join(sentences[:2])  # Use first two sentences as the topic
    return topic

# Step 4: Create PowerPoint presentation with adjusted formatting
def create_presentation(topic, summaries):
    prs = Presentation()

    # Add Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = "Lecture Topic:"
    subtitle_placeholder.text = topic

    # Add each summary as a new slide with two bullet points per slide
    for summary in summaries:
        bullet_points = summary.split('. ')
        for i in range(0, len(bullet_points), 2):  # Change to 3 for three points per slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            body_placeholder = slide.shapes.placeholders[1]

            if i == 0:
                title_placeholder.text = "Key Points"
            body_placeholder.text = "\n".join([f"{point.strip()}" for point in bullet_points[i:i+2]])  # Change to i:i+3 for 3 points per slide

    # Add a Summary slide at the end
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.shapes.placeholders[1]
    title_placeholder.text = "Summary"
    body_placeholder.text = summaries[-1]  # Display the last summary in full

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# Main function to process PDF, summarize and generate PowerPoint
def main(pdf_path):
    pdf_text = extract_text_from_pdf(pdf_path)
    
    # Summarize the text into sections
    summaries = summarize_text(pdf_text)
    
    # Extract the topic from the text
    topic = extract_topic(pdf_text)
    
    # Create the PowerPoint presentation
    ppt_stream = create_presentation(topic, summaries)
    
    # Save to file or return the ppt stream as response
    with open("Marketing.pptx", "wb") as f:
        f.write(ppt_stream.read())

    print("Presentation created successfully! Check the file 'hashing.pptx'.")

# Run the main function
pdf_path = 'Marketing.pdf'  # Provide path to the uploaded PDF
main(pdf_path)